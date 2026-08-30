#!/usr/bin/env python3
"""BC2 Output-Präsentations-Template (KIsult-Stil, Platzhalter {{...}} fuer BC2-Befuellung)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

INK   = RGBColor(0x1A,0x1A,0x2E)
DEEP  = RGBColor(0x10,0x2A,0x43)
BLUE  = RGBColor(0x0F,0x4C,0x81)
ACCENT= RGBColor(0x2E,0x86,0xDE)
TEAL  = RGBColor(0x12,0x9C,0x9C)
GREEN = RGBColor(0x1E,0x88,0x5E)
AMBER = RGBColor(0xE6,0x7E,0x22)
RED   = RGBColor(0xC0,0x39,0x2B)
GREY  = RGBColor(0x5A,0x6B,0x7B)
LIGHT = RGBColor(0xEC,0xF2,0xF9)
CARD  = RGBColor(0xF5,0xF7,0xFA)
WHITE = RGBColor(0xFF,0xFF,0xFF)
LINE  = RGBColor(0x9F,0xB3,0xC8)
PH    = RGBColor(0xFF,0xF6,0xE6)   # placeholder fill
PHL   = RGBColor(0xE6,0x9A,0x2E)   # placeholder border

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]; EMU_W = Inches(13.333)


def slide(): return prs.slides.add_slide(BLANK)

def box(s,x,y,w,h,text,fill,font=WHITE,size=12,bold=True,shape=MSO_SHAPE.ROUNDED_RECTANGLE,
        align=PP_ALIGN.CENTER,line=None,lw=0.75):
    sp=s.shapes.add_shape(shape,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(lw)
    sp.shadow.inherit=False
    tf=sp.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    tf.margin_left=Inches(0.08); tf.margin_right=Inches(0.08); tf.margin_top=Inches(0.03); tf.margin_bottom=Inches(0.03)
    for i,ln in enumerate(text.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        r=p.add_run(); r.text=ln; r.font.size=Pt(size if i==0 else size-2)
        r.font.bold=bold if i==0 else False; r.font.color.rgb=font; r.font.name="Calibri"
    return sp

def txt(s,x,y,w,h,text,size=12,color=INK,bold=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,ln in enumerate(text.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        r=p.add_run(); r.text=ln; r.font.size=Pt(size); r.font.bold=bold
        r.font.color.rgb=color; r.font.name="Calibri"
    return tb

def ph(s,x,y,w,h,text,size=12,align=PP_ALIGN.LEFT):
    """placeholder block"""
    return box(s,x,y,w,h,text,PH,INK,size,False,align=align,line=PHL,lw=1.0)

def header(s,kicker,title,color=BLUE):
    bar=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,EMU_W,Inches(1.1)); bar.fill.solid()
    bar.fill.fore_color.rgb=color; bar.line.fill.background(); bar.shadow.inherit=False
    acc=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,Inches(1.1),EMU_W,Inches(0.05)); acc.fill.solid()
    acc.fill.fore_color.rgb=ACCENT; acc.line.fill.background(); acc.shadow.inherit=False
    txt(s,Inches(0.55),Inches(0.16),Inches(12),Inches(0.3),kicker,12,RGBColor(0xBF,0xD9,0xF2),True)
    txt(s,Inches(0.55),Inches(0.44),Inches(12.2),Inches(0.55),title,24,WHITE,True)

def foot(s,n):
    txt(s,Inches(0.55),Inches(7.08),Inches(9),Inches(0.3),
        "BC2 — Strategic Advisor · Automatisierungs-Workshop  |  {{Kunde}}",9,GREY)
    txt(s,Inches(12.3),Inches(7.08),Inches(0.7),Inches(0.3),str(n),9,GREY,align=PP_ALIGN.RIGHT)

def divider(num,title,sub):
    s=slide()
    bg=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,EMU_W,Inches(7.5)); bg.fill.solid()
    bg.fill.fore_color.rgb=DEEP; bg.line.fill.background(); bg.shadow.inherit=False
    band=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0.9),Inches(3.7),Inches(2.2),Inches(0.09))
    band.fill.solid(); band.fill.fore_color.rgb=ACCENT; band.line.fill.background(); band.shadow.inherit=False
    txt(s,Inches(0.9),Inches(2.4),Inches(4),Inches(1.4),num,90,RGBColor(0x2E,0x6B,0xA8),True)
    txt(s,Inches(0.9),Inches(4.0),Inches(11),Inches(0.9),title,38,WHITE,True)
    txt(s,Inches(0.95),Inches(5.0),Inches(11),Inches(0.6),sub,16,RGBColor(0x9FB9D6 if False else 0x9F,0xB9,0xD6))
    return s

N=[0]
def nxt(): N[0]+=1; return N[0]

# ---------- 1 TITLE ----------
s=slide()
bg=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,EMU_W,Inches(7.5)); bg.fill.solid()
bg.fill.fore_color.rgb=DEEP; bg.line.fill.background(); bg.shadow.inherit=False
band=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,Inches(4.5),EMU_W,Inches(0.08)); band.fill.solid()
band.fill.fore_color.rgb=ACCENT; band.line.fill.background(); band.shadow.inherit=False
txt(s,Inches(0.9),Inches(1.3),Inches(11.5),Inches(0.5),"KI- & AUTOMATISIERUNGS-WORKSHOP",15,RGBColor(0x8F,0xB7,0xE0),True)
txt(s,Inches(0.9),Inches(1.95),Inches(11.5),Inches(1.0),"Zusammenfassung & Potenzialanalyse",40,WHITE,True)
txt(s,Inches(0.9),Inches(3.15),Inches(11.5),Inches(0.6),
    "Ausgangslage · Automatisierungspotenziale · Kostenschätzung",20,RGBColor(0x6E,0xC6,0xE6),True)
ph(s,Inches(0.9),Inches(4.8),Inches(11.5),Inches(1.7),
   "Kunde:        {{Kunde / Unternehmen}}\n"
   "Bereiche:     {{betroffene Bereiche / Prozesse}}\n"
   "Datum:        {{Datum}}        Erstellt von: BC2 — Strategic Advisor (Autonomous CoE Factory)",13)

# ---------- 2 AGENDA ----------
s=slide(); header(s,"AGENDA","Inhalt des Workshops")
items=[("01","Ausgangslage","Zusammenfassung der aktuellen Situation: Unternehmen, Systeme, Herausforderungen",BLUE),
       ("02","Automatisierungspotenziale","Erkannte Potenziale, Beschreibung & Impact, Priorisierung (Komplexität × Impact)",GREEN),
       ("03","Kostenschätzung","Investitionsrahmen je Potenzial und erwarteter Nutzen / Amortisation",AMBER)]
y=Inches(1.6)
for n,t,d,c in items:
    box(s,Inches(0.6),y,Inches(1.2),Inches(1.2),n,c,size=30)
    box(s,Inches(1.95),y,Inches(10.8),Inches(1.2),"",WHITE,INK,11,False,line=LINE)
    txt(s,Inches(2.2),y+Inches(0.16),Inches(10.3),Inches(0.4),t,17,c,True)
    txt(s,Inches(2.2),y+Inches(0.62),Inches(10.3),Inches(0.5),d,12.5,INK)
    y+=Inches(1.55)
foot(s,nxt())

# ---------- 3 DIVIDER 01 ----------
divider("01","Ausgangslage","Zusammenfassung der aktuellen Situation"); nxt()

# ---------- 4 AUSGANGSLAGE FACTS ----------
s=slide(); header(s,"01 · AUSGANGSLAGE","Das Unternehmen — Daten & Fakten")
blocks=[("UNTERNEHMEN & TEAM","{{Branche, Größe, Mitarbeitende,\nTeamstruktur}}",BLUE),
        ("KUNDEN & VOLUMEN","{{Mandanten/Kunden, Fallzahlen,\ntypische Größenordnungen}}",TEAL),
        ("MARKT & ZIELE","{{Marktumfeld, Wettbewerb,\nstrategische Ziele}}",GREEN),
        ("SYSTEMLANDSCHAFT","{{Kernsystem, angebundene Tools,\nSchnittstellen, Reifegrad}}",AMBER)]
xs=[Inches(0.55),Inches(6.75)]; ys=[Inches(1.5),Inches(3.9)]
for i,(t,d,c) in enumerate(blocks):
    x=xs[i%2]; yv=ys[i//2]
    box(s,x,yv,Inches(6.0),Inches(2.25),"",WHITE,INK,11,False,line=c,lw=1.5)
    box(s,x,yv,Inches(6.0),Inches(0.55),t,c,size=13)
    ph(s,x+Inches(0.2),yv+Inches(0.75),Inches(5.6),Inches(1.35),d,12)
foot(s,nxt())

# ---------- 5 HERAUSFORDERUNGEN + KAUSALKETTE ----------
s=slide(); header(s,"01 · AUSGANGSLAGE","Zentrale Herausforderungen & Kausalkette")
txt(s,Inches(0.55),Inches(1.3),Inches(6),Inches(0.35),"Zentrale Herausforderungen",15,BLUE,True)
for i in range(4):
    ph(s,Inches(0.55),Inches(1.75)+Inches(0.72)*i,Inches(6.0),Inches(0.6),
       "{{Herausforderung "+str(i+1)+" — kurz mit Kennzahl}}",12)
txt(s,Inches(6.95),Inches(1.3),Inches(6),Inches(0.35),"Kausalkette: Last → verschenkter Ertrag",15,RED,True)
chain=["URSACHEN","ENGPASS","SINKENDE WIRKUNG","VERSCHENKTER ERTRAG / RISIKO"]
cols=[GREY,AMBER,RED,DEEP]
for i,(t,c) in enumerate(zip(chain,cols)):
    box(s,Inches(6.95),Inches(1.75)+Inches(0.72)*i,Inches(5.85),Inches(0.6),
        t+":  {{...}}",c,size=11.5,align=PP_ALIGN.LEFT)
box(s,Inches(0.55),Inches(4.75),Inches(12.25),Inches(1.55),
    "ANSATZPUNKT / KERN-ZUSAMMENFASSUNG:\n{{Ein bis zwei Sätze: größter Hebel, worauf BC2 die Empfehlung stützt.}}",
    LIGHT,INK,13,False,align=PP_ALIGN.LEFT,line=ACCENT)
foot(s,nxt())

# ---------- 6 DIVIDER 02 ----------
divider("02","Automatisierungspotenziale","Erkennung, Bewertung & Priorisierung"); nxt()

# ---------- 7 POTENZIAL-ÜBERBLICK ----------
s=slide(); header(s,"02 · POTENZIALE","Überblick: erkannte Potenziale",color=GREEN)
txt(s,Inches(0.55),Inches(1.3),Inches(12),Inches(0.35),
    "Von BC2 aus dem Prozessprofil erkannt — pro Kachel ein Potenzial:",13,INK)
for i in range(6):
    col=i%3; row=i//3
    x=Inches(0.55)+Inches(4.18)*col; y=Inches(1.85)+Inches(1.55)*row
    box(s,x,y,Inches(4.0),Inches(1.4),"",WHITE,INK,11,False,line=GREEN,lw=1.25)
    box(s,x+Inches(0.18),y+Inches(0.18),Inches(0.55),Inches(0.55),str(i+1),GREEN,size=16,shape=MSO_SHAPE.OVAL)
    txt(s,x+Inches(0.85),y+Inches(0.2),Inches(3.0),Inches(0.5),"{{Potenzial "+str(i+1)+"}}",13,GREEN,True)
    txt(s,x+Inches(0.2),y+Inches(0.78),Inches(3.6),Inches(0.5),"{{Kurzbeschreibung}}",11,INK)
txt(s,Inches(0.55),Inches(5.1),Inches(12.2),Inches(0.4),
    "Leitprinzip: 80/20 + Mensch-im-Prozess — massive Beschleunigung bei voller Kontrolle.  (Kacheln nach Bedarf duplizieren)",
    11,GREY)
foot(s,nxt())

# ---------- 8 PRIORISIERUNGSMATRIX ----------
s=slide(); header(s,"02 · PRIORISIERUNG","Priorisierungsmatrix: Komplexität × Impact",color=GREEN)
# table
txt(s,Inches(0.55),Inches(1.3),Inches(6),Inches(0.3),"Bewertung je Potenzial (sortiert nach Priorität)",13,BLUE,True)
heads=["Prio","Potenzial","Kompl. 1-10","Impact 1-10"]; hw=[Inches(0.8),Inches(3.0),Inches(1.3),Inches(1.3)]
hx=[Inches(0.55),Inches(1.4),Inches(4.45),Inches(5.8)]
for t,x,w in zip(heads,hx,hw): box(s,x,Inches(1.7),w,Inches(0.42),t,GREEN,size=10.5)
for i in range(6):
    ry=Inches(2.16)+Inches(0.46)*i; fill=WHITE if i%2==0 else LIGHT
    box(s,hx[0],ry,hw[0],Inches(0.44),"{{#}}",fill,INK,10,True,line=LINE)
    box(s,hx[1],ry,hw[1],Inches(0.44),"{{Potenzial}}",fill,INK,10,False,line=LINE,align=PP_ALIGN.LEFT)
    box(s,hx[2],ry,hw[2],Inches(0.44),"{{K}}",fill,INK,10,False,line=LINE)
    box(s,hx[3],ry,hw[3],Inches(0.44),"{{I}}",fill,INK,10,False,line=LINE)
# matrix
mx=Inches(7.4); my=Inches(1.7); mw=Inches(5.3); mh=Inches(4.6)
box(s,mx,my,mw,mh,"",CARD,INK,10,False,line=LINE)
txt(s,mx,my-Inches(0.0)+Inches(0.05),mw,Inches(0.3),"  hoher Impact ↑",10,GREY,True)
txt(s,mx,my+mh-Inches(0.32),mw,Inches(0.3),"  Umsetzungskomplexität  (gering → hoch) →",10,GREY)
box(s,mx+Inches(0.3),my+Inches(0.5),Inches(2.2),Inches(1.7),"★ ZUERST\nQuick Wins\n{{Potenziale}}",GREEN,size=12)
box(s,mx+Inches(2.7),my+Inches(0.5),Inches(2.2),Inches(1.7),"Strategisch\n{{Potenziale}}",BLUE,size=12)
box(s,mx+Inches(0.3),my+Inches(2.4),Inches(2.2),Inches(1.6),"Geringer Hebel\n{{Potenziale}}",GREY,size=12)
box(s,mx+Inches(2.7),my+Inches(2.4),Inches(2.2),Inches(1.6),"Long Bet\n{{Potenziale}}",AMBER,size=12)
foot(s,nxt())

# ---------- 9 POTENZIAL-DETAIL (Vorlage) ----------
s=slide(); header(s,"02 · POTENZIAL · BEWERTUNG","{{Nr}}  {{Titel des Potenzials}}",color=GREEN)
box(s,Inches(11.0),Inches(0.3),Inches(2.0),Inches(0.5),"{{Phase}}",AMBER,size=11)
txt(s,Inches(0.55),Inches(1.35),Inches(6),Inches(0.3),"HEUTIGER PROZESS",12,BLUE,True)
ph(s,Inches(0.55),Inches(1.7),Inches(6.1),Inches(1.85),"{{Wie läuft es heute? Aufwand, Fehlerquote,\nFrequenz — aus BC1-Prozessprofil}}",12)
txt(s,Inches(0.55),Inches(3.75),Inches(6),Inches(0.3),"AUTOMATISIERUNGSPOTENZIAL",12,GREEN,True)
ph(s,Inches(0.55),Inches(4.1),Inches(6.1),Inches(2.0),"{{Was wird automatisiert, in welchen Stufen,\nmit welchem Ergebnis. Human-in-the-Loop.}}",12)
# rating bars
rs=[("MANUELLER AUFWAND HEUTE",BLUE),("EINSPARPOTENZIAL / IMPACT",GREEN),("UMSETZUNGSKOMPLEXITÄT",AMBER)]
for i,(t,c) in enumerate(rs):
    y=Inches(1.7)+Inches(0.95)*i
    box(s,Inches(6.95),y,Inches(5.85),Inches(0.78),"",WHITE,INK,10,False,line=c,lw=1.25)
    txt(s,Inches(7.15),y+Inches(0.08),Inches(4.0),Inches(0.35),t,11,c,True)
    txt(s,Inches(7.15),y+Inches(0.42),Inches(5.5),Inches(0.3),"{{gering | mittel | hoch | sehr hoch}}",12,INK,True)
txt(s,Inches(6.95),Inches(4.6),Inches(6),Inches(0.3),"VORAUSSETZUNGEN / POTENZIELLE LÖSUNG",12,GREY,True)
ph(s,Inches(6.95),Inches(4.95),Inches(5.85),Inches(1.15),"{{Was ist zu klären? Empfohlener Lösungsweg.}}",12)
txt(s,Inches(0.55),Inches(6.25),Inches(12),Inches(0.3),
    "↻ Diese Detail-Folie pro erkanntem Potenzial duplizieren.",11,GREY,True)
foot(s,nxt())

# ---------- 10 DIVIDER 03 ----------
divider("03","Kostenschätzung","Investitionsrahmen & erwarteter Nutzen"); nxt()

# ---------- 11 INVESTITIONSLOGIK ----------
s=slide(); header(s,"03 · KOSTENSCHÄTZUNG","Investitionslogik & erwarteter Nutzen",color=AMBER)
cols=[("KOSTENLOGIK","{{Kalkulationsbasis (z. B. Tagessatz),\nFestpreis je Maßnahme, kleines vs.\ngroßes LLM}}",BLUE),
      ("LAUFENDE KOSTEN","{{LLM-/Token-Nutzung,\nInfrastruktur, Hosting}}",TEAL),
      ("ERWARTETER NUTZEN","{{Zeitersparnis, Fehlerreduktion,\nhöherer Ertrag — je Potenzial}}",GREEN),
      ("DER EIGENTLICHE HEBEL","{{Skalierung & Zukunftssicherheit\nüber reine Zeitersparnis hinaus}}",AMBER)]
xs=[Inches(0.55),Inches(6.75)]; ys=[Inches(1.45),Inches(3.85)]
for i,(t,d,c) in enumerate(cols):
    x=xs[i%2]; yv=ys[i//2]
    box(s,x,yv,Inches(6.0),Inches(2.25),"",WHITE,INK,11,False,line=c,lw=1.5)
    box(s,x,yv,Inches(6.0),Inches(0.55),t,c,size=13)
    ph(s,x+Inches(0.2),yv+Inches(0.75),Inches(5.6),Inches(1.35),d,12)
foot(s,nxt())

# ---------- 12 INVESTITION JE POTENZIAL ----------
s=slide(); header(s,"03 · KOSTENSCHÄTZUNG","Investition je Potenzial (Überblick)",color=AMBER)
heads=["Nr","Potenzial","Phase","Investition (Richtwert)","Erwarteter Nutzen"]
hw=[Inches(0.8),Inches(3.4),Inches(1.5),Inches(2.8),Inches(3.75)]
hx=[Inches(0.55),Inches(1.4),Inches(4.85),Inches(6.4),Inches(9.25)]
for t,x,w in zip(heads,hx,hw): box(s,x,Inches(1.5),w,Inches(0.5),t,AMBER,size=11)
for i in range(7):
    ry=Inches(2.02)+Inches(0.56)*i; fill=WHITE if i%2==0 else LIGHT
    box(s,hx[0],ry,hw[0],Inches(0.54),"{{#}}",fill,INK,10.5,True,line=LINE)
    box(s,hx[1],ry,hw[1],Inches(0.54),"{{Potenzial}}",fill,INK,10.5,False,line=LINE,align=PP_ALIGN.LEFT)
    box(s,hx[2],ry,hw[2],Inches(0.54),"{{Phase}}",fill,INK,10.5,False,line=LINE)
    box(s,hx[3],ry,hw[3],Inches(0.54),"{{€ von–bis}}",fill,INK,10.5,False,line=LINE)
    box(s,hx[4],ry,hw[4],Inches(0.54),"{{Nutzen}}",fill,INK,10.5,False,line=LINE,align=PP_ALIGN.LEFT)
txt(s,Inches(0.55),Inches(6.15),Inches(12.2),Inches(0.5),
    "* Richtwerte/Schätzungen auf Kalkulationsbasis — konkrete Festpreise nach Detaillierung. Werte von BC2 aus Value-Berechnung.",
    10,GREY)
foot(s,nxt())

# ---------- 13 CLOSING ----------
s=slide()
bg=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,EMU_W,Inches(7.5)); bg.fill.solid()
bg.fill.fore_color.rgb=DEEP; bg.line.fill.background(); bg.shadow.inherit=False
band=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0.9),Inches(3.5),Inches(2.2),Inches(0.09))
band.fill.solid(); band.fill.fore_color.rgb=ACCENT; band.line.fill.background(); band.shadow.inherit=False
txt(s,Inches(0.9),Inches(2.4),Inches(11.5),Inches(1.0),"Vielen Dank",44,WHITE,True)
txt(s,Inches(0.9),Inches(3.8),Inches(11.5),Inches(1.6),
    "{{Kern-Botschaft: mehr Ertrag/Effizienz durch erkannte Automatisierungspotenziale —\n"
    "bei voller Kontrolle (Human-in-the-Loop) auf zukunftssicherer Architektur.}}",
    16,RGBColor(0xCF,0xDD,0xEC))
txt(s,Inches(0.9),Inches(6.4),Inches(11.5),Inches(0.4),
    "BC2 — Strategic Advisor · Autonomous CoE Factory",12,RGBColor(0x8F,0xB7,0xE0))

out="/Users/sergio.morazan/Desktop/Privat/Persönliche Entwicklung/Weiterbildung/Studium KI/Architektur BC2/BC2_Praesentations_Template.pptx"
prs.save(out); print("OK ->",out); print("Slides:",len(prs.slides._sldIdLst))
