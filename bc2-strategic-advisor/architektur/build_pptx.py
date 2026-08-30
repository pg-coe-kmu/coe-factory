#!/usr/bin/env python3
"""Erzeugt BC2_Systemarchitektur.pptx (16:9) — v2 inkl. Potenzial-Shift & Erklärungen."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ---------- Palette ----------
INK    = RGBColor(0x1A, 0x1A, 0x2E)
DEEP   = RGBColor(0x16, 0x21, 0x3E)
BLUE   = RGBColor(0x0F, 0x4C, 0x81)
ACCENT = RGBColor(0x2E, 0x86, 0xDE)
TEAL   = RGBColor(0x12, 0x9C, 0x9C)
GREEN  = RGBColor(0x1E, 0x88, 0x5E)
AMBER  = RGBColor(0xE6, 0x7E, 0x22)
RED    = RGBColor(0xC0, 0x39, 0x2B)
PURPLE = RGBColor(0x6C, 0x3A, 0xD6)
GREY   = RGBColor(0x5A, 0x6B, 0x7B)
LIGHT  = RGBColor(0xEC, 0xF2, 0xF9)
CARD   = RGBColor(0xF5, 0xF7, 0xFA)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LINE   = RGBColor(0x9F, 0xB3, 0xC8)
OKGRN  = RGBColor(0xE4, 0xF3, 0xEA)
NOPE   = RGBColor(0xFB, 0xEA, 0xE8)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
EMU_W = Inches(13.333)


def slide():
    return prs.slides.add_slide(BLANK)


def box(s, x, y, w, h, text, fill, font=WHITE, size=12, bold=True,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, align=PP_ALIGN.CENTER, line=None, line_w=0.75):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    tf = sp.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for m in ("margin_left", "margin_right"):
        setattr(tf, m, Inches(0.06))
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size if i == 0 else size - 2)
        r.font.bold = bold if i == 0 else False
        r.font.color.rgb = font; r.font.name = "Calibri"
    return sp


def txt(s, x, y, w, h, text, size=12, color=INK, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = "Calibri"
    return tb


def connect(s, x1, y1, x2, y2, color=GREY, w=1.5, dash=False, arrow=True):
    cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    cn.line.color.rgb = color; cn.line.width = Pt(w); cn.shadow.inherit = False
    ln = cn.line._get_or_add_ln()
    if dash:
        ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    if arrow:
        ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return cn


def header(s, kicker, title, color=BLUE):
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, Inches(1.15))
    bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background(); bar.shadow.inherit = False
    acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.15), EMU_W, Inches(0.06))
    acc.fill.solid(); acc.fill.fore_color.rgb = ACCENT; acc.line.fill.background(); acc.shadow.inherit = False
    txt(s, Inches(0.55), Inches(0.16), Inches(12), Inches(0.3), kicker, 12, RGBColor(0xBF,0xD9,0xF2), True)
    txt(s, Inches(0.55), Inches(0.42), Inches(12.2), Inches(0.6), title, 25, WHITE, True)


def footer(s, n):
    txt(s, Inches(0.55), Inches(7.06), Inches(8), Inches(0.3),
        "BC2 — Strategic Advisor · Systemarchitektur v2", 9, GREY)
    txt(s, Inches(12.3), Inches(7.06), Inches(0.7), Inches(0.3), str(n), 9, GREY, align=PP_ALIGN.RIGHT)


N = [0]
def nxt():
    N[0] += 1
    return N[0]

# ============================================================ 1 — TITLE
s = slide()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, Inches(7.5))
bg.fill.solid(); bg.fill.fore_color.rgb = DEEP; bg.line.fill.background(); bg.shadow.inherit = False
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.55), EMU_W, Inches(0.08))
band.fill.solid(); band.fill.fore_color.rgb = ACCENT; band.line.fill.background(); band.shadow.inherit = False
txt(s, Inches(0.9), Inches(1.4), Inches(11.5), Inches(0.5),
    "AUTONOMOUS CoE FACTORY · PG KI-CoE-KMU", 15, RGBColor(0x8F,0xB7,0xE0), True)
txt(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.3), "BC2 — Strategic Advisor", 48, WHITE, True)
txt(s, Inches(0.9), Inches(3.2), Inches(11.5), Inches(0.8), "Systemarchitektur v2", 28, RGBColor(0x6E,0xC6,0xE6), True)
txt(s, Inches(0.9), Inches(4.8), Inches(11.6), Inches(1.7),
    "Ziel: Automatisierungspotenziale erkennen & Value berechnen → Präsentation\n"
    "Stack: n8n · Claude (Sonnet 4.6 / Opus 4.7) · PostgreSQL · FastAPI\n"
    "Team: Eike Bischof · Sergio Morazán Irias   |   Stand: 27.06.2026",
    14, RGBColor(0xCF,0xDD,0xEC))

# ============================================================ 2 — ZIELBILD / OUTPUT
s = slide()
header(s, "0 · ZIELBILD", "Was BC2 produziert: eine entscheidungsreife Präsentation", color=GREEN)
txt(s, Inches(0.55), Inches(1.4), Inches(12.2), Inches(0.5),
    "BC2 nimmt das Prozessprofil aus BC1 und erzeugt ein KIsult-artiges Workshop-Deck mit drei Teilen:",
    14, INK)
cards = [
    ("1", "AUSGANGSLAGE", "Zusammenfassung der aktuellen Situation:\nUnternehmen, Systeme, Prozesse,\nHerausforderungen", BLUE),
    ("2", "AUTOMATISIERUNGS-\nPOTENZIALE", "Erkannte Potenziale je mit\nBeschreibung + IMPACT, Aufwand-heute,\nKomplexität, Priorisierung", GREEN),
    ("3", "KOSTENSCHÄTZUNG", "Investitionsrahmen je Potenzial\n(Richtwert) + erwarteter Nutzen /\nAmortisation", AMBER),
]
x = Inches(0.55)
for n, t, d, c in cards:
    box(s, x, Inches(2.15), Inches(3.95), Inches(2.7), "", WHITE, INK, 11, False, line=c, line_w=1.5)
    box(s, x + Inches(0.25), Inches(2.4), Inches(0.7), Inches(0.7), n, c, size=22, shape=MSO_SHAPE.OVAL)
    txt(s, x + Inches(1.1), Inches(2.45), Inches(2.7), Inches(0.9), t, 15, c, True, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x + Inches(0.3), Inches(3.45), Inches(3.4), Inches(1.3), d, 12, INK)
    x += Inches(4.1)
box(s, Inches(0.55), Inches(5.2), Inches(12.2), Inches(1.45),
    "KEIN fester Musterkatalog mehr:  BC2 ordnet nicht vorgegebene Automatisierungs-MUSTER zu, sondern\n"
    "ERKENNT Potenziale aus dem konkreten Prozess und BERECHNET deren möglichen Value.\n"
    "→ ersetzt die ursprüngliche Pattern-Matching-Engine + Qdrant.   Parallel: maschinenlesbares konzept.json an BC3.",
    OKGRN, INK, 13, False, align=PP_ALIGN.LEFT, line=GREEN)
footer(s, nxt())

# ============================================================ 3 — CONTEXT
s = slide()
header(s, "1 · CONTEXT VIEW", "Einordnung — und wohin Gate 0 gehört")
mid = Inches(3.1)
box(s, Inches(0.55), Inches(2.6), Inches(2.0), Inches(1.0), "Process Owner\n/ BC1", GREY, size=13)
box(s, Inches(2.95), Inches(2.4), Inches(1.8), Inches(1.4), "Gate 0\nDatenqualität\n→ BC1/Platform", AMBER, size=11, shape=MSO_SHAPE.HEXAGON)
box(s, Inches(5.2), Inches(2.6), Inches(2.7), Inches(1.0), "BC2\nStrategic Advisor", BLUE, size=16)
box(s, Inches(8.35), Inches(2.4), Inches(1.7), Inches(1.4), "Gate 1\nBusiness\nApproval", PURPLE, size=11, shape=MSO_SHAPE.HEXAGON)
box(s, Inches(10.5), Inches(2.6), Inches(2.3), Inches(1.0), "BC3\nEngineering\nArchitect", GREEN, size=14)
box(s, Inches(8.5), Inches(4.4), Inches(1.4), Inches(0.7), "CoE-Manager\n(HitL)", TEAL, size=11, shape=MSO_SHAPE.OVAL)
connect(s, Inches(2.55), mid, Inches(2.95), mid, ACCENT, 2)
connect(s, Inches(4.75), mid, Inches(5.2), mid, ACCENT, 2)
connect(s, Inches(7.9), mid, Inches(8.35), mid, ACCENT, 2)
connect(s, Inches(10.05), mid, Inches(10.5), mid, ACCENT, 2)
connect(s, Inches(9.2), Inches(4.4), Inches(9.2), Inches(3.8), TEAL, 1.5, dash=True)
txt(s, Inches(2.9), Inches(3.85), Inches(2.0), Inches(0.3), "bc1.profile.ready", 9, GREY)
txt(s, Inches(7.6), Inches(3.85), Inches(2.6), Inches(0.3), "konzept.json + Präsentation", 9, GREY)
box(s, Inches(0.55), Inches(5.35), Inches(12.2), Inches(1.3),
    "KORREKTUR — Gate 0 gehört zu BC1/Platform, nicht zu BC2:\n"
    "Der Produzent (BC1) garantiert die Vertragsqualität und feuert bc1.profile.ready nur, wenn die Qualität stimmt.\n"
    "BC2 macht an seiner Grenze nur eine schlanke, defensive Eingangsvalidierung (Schema + Pflichtfelder).",
    LIGHT, INK, 13, False, align=PP_ALIGN.LEFT, line=AMBER)
footer(s, nxt())

# ============================================================ 4 — COMPONENTS
s = slide()
header(s, "2 · COMPONENT VIEW", "Komponentenarchitektur innerhalb BC2")

def layer(x, y, w, h, label, col):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = RGBColor(0xEE,0xF3,0xF8)
    sp.line.color.rgb = LINE; sp.line.width = Pt(1); sp.shadow.inherit = False
    tb = sp.text_frame; tb.vertical_anchor = MSO_ANCHOR.TOP
    p = tb.paragraphs[0]; r = p.add_run(); r.text = label
    r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = col; r.font.name = "Calibri"

cy = Inches(1.55)
layer(Inches(0.4), cy, Inches(3.0), Inches(2.5), "API & EINGANG (FastAPI)", BLUE)
box(s, Inches(0.55), Inches(1.95), Inches(2.7), Inches(0.6), "REST-Endpunkte", BLUE, size=11)
box(s, Inches(0.55), Inches(2.62), Inches(2.7), Inches(0.6), "Event-Sub · bc1.profile.ready", ACCENT, size=10)
box(s, Inches(0.55), Inches(3.3), Inches(2.7), Inches(0.6), "Input-Validator (defensiv)", AMBER, size=10)

layer(Inches(3.6), cy, Inches(2.4), Inches(2.5), "ORCHESTRIERUNG (n8n)", TEAL)
box(s, Inches(3.75), Inches(1.95), Inches(2.1), Inches(1.95),
    "n8n Workflow\nbc2-advisor\n\nLoad → Analyse →\nPotenziale → Value →\nPrio → Präsentation", TEAL, size=10.5)

layer(Inches(6.2), cy, Inches(3.85), Inches(2.5), "FACHLOGIK (Python 3.11)", GREEN)
fl = [("Situations-Analyse", BLUE), ("Potenzial-Erkennung", GREEN),
      ("Value-/ROI-Berechnung", GREEN), ("Priorisierung", TEAL),
      ("Präsentations-Generator", PURPLE), ("Schema-Validator", GREY)]
for i, (t, c) in enumerate(fl):
    col = i % 2; row = i // 2
    box(s, Inches(6.35) + Inches(1.78)*col, Inches(1.95) + Inches(0.66)*row,
        Inches(1.68), Inches(0.56), t, c, size=9.5)

layer(Inches(10.25), cy, Inches(2.65), Inches(2.5), "LLM-LAYER", PURPLE)
box(s, Inches(10.4), Inches(1.95), Inches(2.35), Inches(1.0), "Claude API\nSonnet 4.6 default\nOpus 4.7 Reasoning", PURPLE, size=10.5)
box(s, Inches(10.4), Inches(3.05), Inches(2.35), Inches(0.85), "PostgreSQL\nkonzepte · audit_log", DEEP, size=10.5)

box(s, Inches(0.4), Inches(4.35), Inches(6.1), Inches(0.95),
    "HUMAN-IN-THE-LOOP — Gate-1-Dashboard\nApprove / Reject · Budget · Audit-Log  (BC2s einziges Frontend)", PURPLE, size=12)
box(s, Inches(6.7), Inches(4.35), Inches(6.2), Inches(0.95),
    "OUTPUT\nPräsentation (3 Teile) für CoE-Manager  +  konzept.json → BC3", GREEN, size=12)
connect(s, Inches(3.4), Inches(2.8), Inches(3.6), Inches(2.8), GREY, 1.5)
connect(s, Inches(6.0), Inches(2.8), Inches(6.2), Inches(2.8), GREY, 1.5)
connect(s, Inches(10.05), Inches(2.8), Inches(10.25), Inches(2.8), GREY, 1.5)
txt(s, Inches(0.4), Inches(5.5), Inches(12.5), Inches(0.4),
    "Potenzial-Erkennung ersetzt die alte Pattern-Matching-Engine (kein Qdrant/Musterkatalog).  "
    "Value/ROI rechnet deterministisch in Python; Claude liefert nur qualitative Beschreibung & Impact.",
    11, INK)
footer(s, nxt())

# ============================================================ 5 — PIPELINE
s = slide()
header(s, "3 · RUNTIME VIEW", "Laufzeit-Pipeline (n8n bc2-advisor)", color=TEAL)
steps = [
    ("1", "Load", "Prozessprofil\naus Postgres", BLUE),
    ("2", "Validate", "Schema /\nPflichtfelder", AMBER),
    ("3", "Analyse", "Ist-Situation\n→ Teil 1", BLUE),
    ("4", "Potenziale", "erkennen (LLM)\n→ Teil 2", GREEN),
    ("5", "Value", "ROI / Kosten\ndeterministisch", GREEN),
    ("6", "Prio", "Komplexität\nx Impact", TEAL),
    ("7", "Präsentation", "3 Teile +\nkonzept.json", PURPLE),
    ("8", "Persist", "Postgres\ngate1=pending", BLUE),
]
x0 = Inches(0.45); y0 = Inches(1.65); w = Inches(2.95); h = Inches(1.2); gx = Inches(0.13)
for i, (n, t, d, c) in enumerate(steps):
    col = i % 4; row = i // 4
    bx = x0 + (w + gx) * col; by = y0 + Inches(1.6) * row
    box(s, bx, by, w, h, "", CARD, INK, 10, False, line=LINE)
    box(s, bx, by, Inches(0.42), Inches(0.42), n, c, size=13, shape=MSO_SHAPE.OVAL)
    txt(s, bx + Inches(0.52), by + Inches(0.05), w - Inches(0.6), Inches(0.4), t, 13, c, True)
    txt(s, bx + Inches(0.12), by + Inches(0.52), w - Inches(0.2), Inches(0.6), d, 11, INK)
    if col < 3 and i < len(steps) - 1:
        connect(s, bx + w, by + Inches(0.6), bx + w + gx, by + Inches(0.6), TEAL, 1.6)
connect(s, x0 + (w+gx)*3 + Inches(1.4), y0 + h, x0 + (w+gx)*3 + Inches(1.4), y0 + Inches(1.6), TEAL, 1.6)
box(s, Inches(0.45), Inches(5.15), Inches(6.1), Inches(1.4),
    "GATE 1 · Human-in-the-Loop\nCoE-Manager: Approve / Reject + Budget/Kommentar\n→ approved: gate1=approved, audit_log, bc2.konzept.ready → BC3",
    GREEN, size=12, align=PP_ALIGN.LEFT)
box(s, Inches(6.75), Inches(5.15), Inches(6.1), Inches(1.4),
    "REJECT-LOOP\nKommentar → zurück an Potenzial-/Reasoning-Stufe (4)\nPräsentation wird neu erzeugt, erneut zur Freigabe vorgelegt",
    AMBER, size=12, align=PP_ALIGN.LEFT)
footer(s, nxt())

# ============================================================ 6 — POTENZIAL SHIFT
s = slide()
header(s, "KONZEPT-SHIFT", "Potenziale erkennen statt Muster zuordnen", color=GREEN)
box(s, Inches(0.55), Inches(1.5), Inches(5.9), Inches(2.2), "", NOPE, INK, 11, False, line=RED, line_w=1.5)
txt(s, Inches(0.8), Inches(1.65), Inches(5.4), Inches(0.4), "ALT  ·  Musterkatalog (verworfen)", 15, RED, True)
txt(s, Inches(0.8), Inches(2.15), Inches(5.5), Inches(1.5),
    "Schmerzpunkt → Vektor-Suche in Qdrant\ngegen ~20 fixe Muster (RPA/API/…)\n→ Top-N Treffer.\n\n"
    "Nachteil: starrer Katalog nötig, Qdrant-\nBetrieb, wenig anschlussfähig an ein\nWorkshop-Deck.", 12, INK)
box(s, Inches(6.85), Inches(1.5), Inches(5.95), Inches(2.2), "", OKGRN, INK, 11, False, line=GREEN, line_w=1.5)
txt(s, Inches(7.1), Inches(1.65), Inches(5.4), Inches(0.4), "NEU  ·  Potenzial-Erkennung + Value", 15, GREEN, True)
txt(s, Inches(7.1), Inches(2.15), Inches(5.5), Inches(1.5),
    "Claude analysiert das konkrete Prozess-\nprofil und ERKENNT Potenziale (offen, nicht\naus Liste) + begründet Impact.\n\n"
    "Value/Kosten werden deterministisch in\nPython gerechnet. Kein Qdrant, kein\ngepflegter Katalog.", 12, INK)
connect(s, Inches(6.45), Inches(2.6), Inches(6.85), Inches(2.6), GREEN, 2.5)
txt(s, Inches(0.55), Inches(3.95), Inches(12.2), Inches(0.4),
    "Pro Potenzial bewertet (wie im KIsult-Workshop):", 14, BLUE, True)
chips = [("Manueller Aufwand heute", BLUE), ("Einsparpotenzial / Impact", GREEN),
         ("Umsetzungskomplexität", AMBER), ("Investition (Richtwert)", PURPLE),
         ("Amortisation", TEAL), ("Kategorie: Quick Win / Strategisch / …", GREY)]
x = Inches(0.55); y = Inches(4.5)
for i, (t, c) in enumerate(chips):
    col = i % 3; row = i // 3
    box(s, Inches(0.55) + Inches(4.1)*col, y + Inches(0.78)*row, Inches(3.9), Inches(0.62), t, c, size=12)
txt(s, Inches(0.55), Inches(6.25), Inches(12.2), Inches(0.5),
    "Konsequenz: konzept.schema.json wird Potenzial-zentriert (Feld 'value{}' statt 'empfohlenes_muster').",
    11, INK, True)
footer(s, nxt())

# ============================================================ 7 — API & FASTAPI
s = slide()
header(s, "4 · API-DESIGN", "REST-Endpunkte (FastAPI)", color=BLUE)
box(s, Inches(0.55), Inches(1.45), Inches(12.25), Inches(0.95),
    "Was ist FastAPI?  Ein modernes Python-Framework für REST-APIs: Funktion mit Typ-Hints schreiben → FastAPI erzeugt\n"
    "den HTTP-Endpunkt UND automatisch die OpenAPI-/Swagger-Doku (unter /docs). Leichtgewichtig, async-fähig.",
    LIGHT, INK, 12, False, align=PP_ALIGN.LEFT, line=LINE)
rows = [
    ("POST", "/bc2/konzept", "Konzept-Erzeugung anstoßen → startet Pipeline, gibt konzept_id", GREEN),
    ("GET", "/bc2/konzept/{id}", "Ein Konzept abrufen (Dashboard & BC3)", BLUE),
    ("GET", "/bc2/konzepte?status=pending", "Liste offener Konzepte (fürs Gate 1)", BLUE),
    ("GET", "/bc2/konzept/{id}/praesentation", "Präsentation (3 Teile) — JSON / Markdown / PDF", PURPLE),
    ("GET", "/bc2/priorisierung/{id}", "Priorisierung (L2-02)", TEAL),
    ("POST", "/bc2/konzept/{id}/gate1", "Gate-1-Entscheidung → feuert bc2.konzept.ready", AMBER),
    ("GET", "/healthz · /readyz", "Health-Checks", GREY),
]
y = Inches(2.65)
for i, (m, p, d, c) in enumerate(rows):
    ry = y + Inches(0.6)*i
    box(s, Inches(0.55), ry, Inches(1.1), Inches(0.5), m, c, size=11)
    box(s, Inches(1.75), ry, Inches(4.0), Inches(0.5), p, WHITE, INK, 11, True, line=LINE, align=PP_ALIGN.LEFT)
    txt(s, Inches(5.95), ry + Inches(0.04), Inches(6.9), Inches(0.45), d, 11.5, INK, anchor=MSO_ANCHOR.MIDDLE)
footer(s, nxt())

# ============================================================ 8 — FRONTEND
s = slide()
header(s, "5 · FRONTEND / HitL", "Gibt es ein Frontend für Gate 0 und Gate 1?", color=PURPLE)
box(s, Inches(0.55), Inches(1.6), Inches(5.95), Inches(4.6), "", NOPE, INK, 11, False, line=RED, line_w=1.5)
box(s, Inches(0.55), Inches(1.6), Inches(5.95), Inches(0.8), "GATE 0  —  KEIN Frontend", RED, size=18)
txt(s, Inches(0.85), Inches(2.65), Inches(5.4), Inches(3.4),
    "• Automatischer Schwellwert-Check\n  (Vollständigkeit, PII)\n\n"
    "• Keine menschliche Entscheidung\n  → kein UI nötig\n\n"
    "• Liegt ohnehin bei BC1 / Platform\n\n"
    "• In BC2 nur als defensiver\n  Input-Validator vorhanden", 14, INK)
box(s, Inches(6.85), Inches(1.6), Inches(5.95), Inches(4.6), "", OKGRN, INK, 11, False, line=GREEN, line_w=1.5)
box(s, Inches(6.85), Inches(1.6), Inches(5.95), Inches(0.8), "GATE 1  —  JA, Dashboard", GREEN, size=18)
txt(s, Inches(7.15), Inches(2.65), Inches(5.4), Inches(3.4),
    "• BC2s einziges echtes Frontend\n  (AP 2.6 · L2-04)\n\n"
    "• Web-UI für den CoE-Manager:\n  Potenziale, Impact, ROI je Konzept\n\n"
    "• Approve / Reject + Budget anpassen\n\n"
    "• Schreibt Audit-Log, feuert bei\n  Approve bc2.konzept.ready", 14, INK)
footer(s, nxt())

# ============================================================ 9 — MVP vs ZIEL
s = slide()
header(s, "6 · MVP vs. ZIEL", "Welche Komponente braucht ihr wirklich?", color=DEEP)
txt(s, Inches(0.55), Inches(1.35), Inches(12.2), Inches(0.4),
    "2-Personen-Team, Walking Skeleton bis M2 — vieles ist Ziel-Bild und für den MVP optional.", 13, INK)
hdr = [("Komponente", Inches(2.4)), ("Wofür", Inches(5.7)), ("MVP?", Inches(1.3)), ("Erklärung", Inches(2.6))]
xs = [Inches(0.55), Inches(2.95), Inches(8.65), Inches(9.95)]
yy = Inches(1.85)
for (t, w), x in zip(hdr, xs):
    box(s, x, yy, w, Inches(0.45), t, DEEP, size=11)
data = [
    ("FastAPI", "REST-Endpunkte bereitstellen", "JA", "API-Technik", GREEN),
    ("PostgreSQL", "Konzepte · Priorisierung · audit_log", "JA", "geteilte DB", GREEN),
    ("Claude", "Analyse · Potenziale · Texte", "JA", "Kern-Logik", GREEN),
    ("n8n", "Pipeline-Orchestrierung (dünn)", "optional", "MVP: Python", AMBER),
    ("Redis", "Event-Bus / LLM-Cache", "NEIN", "HTTP reicht", RED),
    ("Qdrant", "Vektor-DB", "NEIN", "entfällt ganz", RED),
    ("Grafana/Loki", "Observability / Kosten", "NEIN", "JSON-Logs", RED),
]
for i, (k, wof, mvp, erk, c) in enumerate(data):
    ry = yy + Inches(0.5) + Inches(0.52)*i
    fill = WHITE if i % 2 == 0 else LIGHT
    box(s, xs[0], ry, Inches(2.4), Inches(0.48), k, fill, INK, 11, True, line=LINE, align=PP_ALIGN.LEFT)
    box(s, xs[1], ry, Inches(5.7), Inches(0.48), wof, fill, INK, 10.5, False, line=LINE, align=PP_ALIGN.LEFT)
    box(s, xs[2], ry, Inches(1.3), Inches(0.48), mvp, c, WHITE, 11, True)
    box(s, xs[3], ry, Inches(2.6), Inches(0.48), erk, fill, INK, 10.5, False, line=LINE, align=PP_ALIGN.LEFT)
footer(s, nxt())

# ============================================================ 10 — TECH-ENTSCHEIDUNGEN
s = slide()
header(s, "7 · TECH-ENTSCHEIDUNGEN", "n8n vs. LangChain — und warum", color=BLUE)
box(s, Inches(0.55), Inches(1.5), Inches(6.0), Inches(2.6), "", WHITE, INK, 11, False, line=TEAL, line_w=1.5)
txt(s, Inches(0.8), Inches(1.65), Inches(5.5), Inches(0.4), "n8n  (Low-Code, visuell)", 15, TEAL, True)
txt(s, Inches(0.8), Inches(2.15), Inches(5.5), Inches(1.9),
    "+ Connectors, Retry, Trigger out-of-the-box\n+ gut für LINEARE Abläufe (genau unser Fall)\n+ Workflow als versioniertes JSON\n"
    "– komplexe Logik wird unübersichtlich\n– Unit-Tests (DoD ≥70 %) kaum möglich", 12.5, INK)
box(s, Inches(6.8), Inches(1.5), Inches(6.0), Inches(2.6), "", WHITE, INK, 11, False, line=PURPLE, line_w=1.5)
txt(s, Inches(7.05), Inches(1.65), Inches(5.5), Inches(0.4), "LangChain / LangGraph  (Code-First)", 15, PURPLE, True)
txt(s, Inches(7.05), Inches(2.15), Inches(5.5), Inches(1.9),
    "+ stark bei komplexen Agenten-Loops\n+ Tool-Use, volle Testbarkeit\n– mehr Boilerplate, kein visuelles Bild\n"
    "– für unseren LINEAREN Ablauf Overkill\n→ erst relevant, wenn BC2 echt agentisch wird", 12.5, INK)
box(s, Inches(0.55), Inches(4.35), Inches(12.25), Inches(1.0),
    "EMPFEHLUNG (Hybrid):  n8n nur als dünner Orchestrator (Trigger, Claude-Call, Retry) — die echte Logik\n"
    "(Potenzial-Erkennung, Value, Validierung) in getesteten Python-Modulen. Würde BC2 später agentisch: LangGraph statt n8n.",
    OKGRN, INK, 13, False, align=PP_ALIGN.LEFT, line=GREEN)
box(s, Inches(0.55), Inches(5.55), Inches(12.25), Inches(1.05),
    "Warum getrennt?  Value/ROI deterministisch in Python = reproduzierbar & testbar (Unit-Tests gegen Mock-Profil).\n"
    "Das LLM (Claude) liefert nur die qualitative Beschreibung & Impact-Einordnung — nicht die Zahlen.",
    LIGHT, INK, 13, False, align=PP_ALIGN.LEFT, line=LINE)
footer(s, nxt())

# ============================================================ 11 — AP MAPPING
s = slide()
header(s, "8 · AP-MAPPING", "Architektur ↔ Arbeitspakete & Lieferung", color=GREEN)
hy = Inches(1.7)
box(s, Inches(0.55), hy, Inches(1.4), Inches(0.5), "AP", GREEN, size=12)
box(s, Inches(2.05), hy, Inches(5.6), Inches(0.5), "Architektur-Baustein", GREEN, size=12)
box(s, Inches(7.75), hy, Inches(5.05), Inches(0.5), "Liefergegenstand", GREEN, size=12)
rows = [
    ("2.1", "Schnittstellen-Contract (Schemata)", "konzept/priorisierung.schema.json"),
    ("2.2", "Potenzial-Erkennung (LLM)", "erkannte Potenziale (Teil 2)"),
    ("2.3", "Value-/ROI-Berechnung", "Kostenschätzung (Teil 3)"),
    ("2.4", "Situations-Analyse + Präsentations-Generator", "Ausgangslage (Teil 1) + Gesamt-Deck"),
    ("2.5", "Priorisierung (Komplexität × Impact)", "L2-02 Prozesspriorisierung"),
    ("2.6", "Gate-1-Dashboard (HitL)", "L2-04 Gate-1-Dashboard"),
    ("2.7", "End-to-End-Pipeline", "verifizierte Übergabe BC1→BC2→BC3"),
    ("2.8", "Präsentations-Output", "L2-05 Ergebnis-Präsentation"),
]
for i, (ap, b, d) in enumerate(rows):
    ry = hy + Inches(0.56) + Inches(0.55)*i
    fill = WHITE if i % 2 == 0 else LIGHT
    box(s, Inches(0.55), ry, Inches(1.4), Inches(0.5), "AP "+ap, fill, INK, 11, True, line=LINE)
    box(s, Inches(2.05), ry, Inches(5.6), Inches(0.5), b, fill, INK, 11, False, line=LINE, align=PP_ALIGN.LEFT)
    box(s, Inches(7.75), ry, Inches(5.05), Inches(0.5), d, fill, INK, 11, False, line=LINE, align=PP_ALIGN.LEFT)
footer(s, nxt())

out = "/Users/sergio.morazan/Desktop/Privat/Persönliche Entwicklung/Weiterbildung/Studium KI/Architektur BC2/BC2_Systemarchitektur.pptx"
prs.save(out)
print("OK ->", out)
print("Slides:", len(prs.slides._sldIdLst))
